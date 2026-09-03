"""The framework deck: other definitions, ours, the pipeline, findings, next steps.

Built for import into Google Slides, so everything is a plain text box, a plain
table or a picture. No native chart objects, no SmartArt, no theme fonts: Google
re-flows those on import and the layout drifts. Charts arrive as PNG, already
rendered by make_framework_charts.py from out/panel_run.

Every figure on a slide traces to a file in that run directory. Nothing is typed
from memory.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

CHARTS = Path(sys.argv[1] if len(sys.argv) > 1 else "reporting/charts/framework")
OUT = Path(sys.argv[2] if len(sys.argv) > 2
           else "out/analysis/dnssec_framework.pptx")

INK = RGBColor(0x17, 0x1A, 0x14)
MUTED = RGBColor(0x66, 0x70, 0x5C)
RULE = RGBColor(0xD3, 0xD8, 0xCC)
PAPER = RGBColor(0xF6, 0xF7, 0xF3)
TEAL = RGBColor(0x00, 0x90, 0x7F)
AMBER = RGBColor(0xC9, 0x6A, 0x06)
BRICK = RGBColor(0xB3, 0x2B, 0x22)
FONT = "Segoe UI"

#: Paragraph break inside a textbox string. Spelled chr(10) rather than an
#: escape so it survives being edited through shells and heredocs, which have
#: mangled it more than once.
NL = chr(10)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height
M = Inches(0.8)
BLANK = prs.slide_layouts[6]


def textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
            align=PP_ALIGN.LEFT, space_after=6, line=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        run = p.add_run()
        run.text = para
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return tb


def slide_title(slide, title, dek=None, eyebrow=None):
    y = M - Inches(0.25)
    if eyebrow:
        textbox(slide, M, y, W - 2 * M, Inches(0.3), eyebrow.upper(),
                size=11, bold=True, color=TEAL)
        y += Inches(0.38)
    textbox(slide, M, y, W - 2 * M, Inches(0.7), title, size=32, bold=True)
    y += Inches(0.82)
    if dek:
        textbox(slide, M, y, W - 2 * M, Inches(0.5), dek, size=15, color=MUTED)
        y += Inches(0.62)
    line = slide.shapes.add_shape(1, M, y, W - 2 * M, Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = RULE
    line.line.fill.background(); line.shadow.inherit = False
    return y + Inches(0.34)


def new(title=None, dek=None, eyebrow=None):
    s = prs.slides.add_slide(BLANK)
    top = slide_title(s, title, dek, eyebrow) if title else M
    return s, top


def picture(slide, name, top, *, max_h=None, max_w=None):
    """Place a chart centred, scaled to fit, never stretched."""
    path = CHARTS / f"{name}.png"
    iw, ih = Image.open(path).size
    avail_w = max_w or (W - 2 * M)
    avail_h = max_h or (H - top - Inches(0.75))
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int((W - w) / 2), int(top), w, h)
    return top + h


def table(slide, x, y, w, rows, widths, *, size=13, header=True, row_h=0.42):
    """A plain pptx table. Google Slides imports these faithfully."""
    n_r, n_c = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_r, n_c, x, y, w, Inches(row_h * n_r))
    tbl = shape.table
    total = sum(widths)
    for i, frac in enumerate(widths):
        tbl.columns[i].width = Emu(int(w * frac / total))
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(row_h)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(size if not (header and r == 0) else size - 1)
            run.font.name = FONT
            run.font.bold = header and r == 0
            run.font.color.rgb = MUTED if (header and r == 0) else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                PAPER if (header and r == 0) else RGBColor(0xFF, 0xFF, 0xFF))
    return y + Inches(row_h * n_r)


def cards(slide, top, items, *, height=1.85, size=14):
    """Equal cards across the width: title, body, optional footer."""
    n = len(items)
    gap = Inches(0.22)
    cw = int((W - 2 * M - gap * (n - 1)) / n)
    for i, (head, body, foot) in enumerate(items):
        x = M + i * (cw + gap)
        box = slide.shapes.add_shape(1, x, top, cw, Inches(height))
        box.fill.solid(); box.fill.fore_color.rgb = PAPER
        box.line.color.rgb = RULE; box.line.width = Pt(0.75)
        box.shadow.inherit = False
        box.text_frame.text = ""
        textbox(slide, x + Inches(0.18), top + Inches(0.14),
                cw - Inches(0.36), Inches(0.3), head, size=12, bold=True,
                color=TEAL)
        textbox(slide, x + Inches(0.18), top + Inches(0.52),
                cw - Inches(0.36), Inches(0.9), body, size=size)
        if foot:
            textbox(slide, x + Inches(0.18), top + Inches(height) - Inches(0.62),
                    cw - Inches(0.36), Inches(0.5), foot, size=11.5, color=MUTED)
    return top + Inches(height)


# =========================================================================== #
# 1. Title
# =========================================================================== #
s, _ = new()
textbox(s, M, Inches(2.1), W - 2 * M, Inches(0.4),
        "DNSSEC RFC ADOPTION — MEASUREMENT FRAMEWORK", size=13, bold=True,
        color=TEAL)
textbox(s, M, Inches(2.6), W - 2 * M, Inches(1.6),
        "Measuring what the data can\nactually tell us", size=44, bold=True)
textbox(s, M, Inches(4.4), Inches(8.2), Inches(1.2),
        "Why the existing definitions do not transfer, the one we built instead, "
        "the pipeline that computes it, and what it has found so far.",
        size=17, color=MUTED)
textbox(s, M, Inches(6.2), Inches(9), Inches(0.5),
        "Full run: 8.24 TB · 15,179 source-days · 25,151 files · 12 sources · "
        "5 RIRs (2009–2026) + 7 forward TLDs (2016–2023)", size=12, color=MUTED)

# =========================================================================== #
# 2. The problem
# =========================================================================== #
s, top = new("One word was doing three jobs",
             "Every sentence with “adoption” in it hid which of these was meant.",
             eyebrow="the problem")
cards(s, top, [
    ("SOMEONE TRIED IT",
     "Ed25519 first appeared in 2019-01, 1.9 years after RFC 8080.",
     "true, and almost meaningless alone"),
    ("IT REACHED REAL USE",
     "In the forward TLDs it reached 2.04% of signed names by 2020-06.",
     "a different measure entirely"),
    ("IT IS SPREADING",
     "It is not. Those same TLDs were back to 0.003% by 2023-12.",
     "a third measure again"),
], height=2.0)
textbox(s, M, top + Inches(2.4), W - 2 * M, Inches(1.4),
        "“Ed25519 was adopted after 1.9 years” and “Ed25519 was never adopted” "
        "are both defensible readings of the same record.\n"
        "That is what makes the word unusable once a number is attached to it.",
        size=17)

# =========================================================================== #
# 3. Other definitions — and how they were derived
# =========================================================================== #
s, top = new("What already exists, and where it comes from",
             "The method behind each definition decides what it can and cannot say.",
             eyebrow="other definitions")
after = table(s, M, top, W - 2 * M, [
    ["Framework", "How it was derived", "What it measures", "Fits us?"],
    ["Rogers / Bass\ndiffusion",
     "Theory + market data on\nconsumer goods",
     "Cumulative adopters —\ncan only go up",
     "No"],
    ["ISA model\nHovav et al. 2004",
     "SURVEYS of organisations\n(Australia, China)",
     "Why managers decide:\nattitude, top-mgmt support",
     "No"],
    ["Osterweil et al.\nIMC 2008",
     "ACTIVE DNS QUERIES\nfrom resolvers",
     "availability / verifiability /\nvalidity — does it work now",
     "Different\ninstrument"],
    ["Chung et al.\nIMC 2017",
     "Measurement + buying\ndomains from registrars",
     "Who is responsible\n(registrar attribution)",
     "Out of\nscope"],
    ["RFC 5218",
     "Retrospective on IETF\nprotocol histories",
     "Why protocols succeed\n(“incrementally deployable”)",
     "Partly —\nwe cite it"],
], [0.17, 0.24, 0.35, 0.14], row_h=0.58)
textbox(s, M, after + Inches(0.14), W - 2 * M, Inches(0.72),
        "The ISA model asks people why they deployed. We read what zones publish. "
        "Nothing in a zone file records intent, management support, or attitude — "
        "so a survey-derived definition cannot be computed from our data at all.",
        size=13.5, color=MUTED)

# =========================================================================== #
# 4. Why the diffusion curve does not transfer
# =========================================================================== #
s, top = new("Their curve can only go up. Ours goes down.",
             "The same mechanism, measured two ways.",
             eyebrow="why they don't transfer")
picture(s, "why_not_rogers", top, max_h=Inches(3.9))
textbox(s, M, H - Inches(1.55), W - 2 * M, Inches(1.1),
        "RSA/SHA-256 went 0% → 74.6% → 22.3%. Nothing failed: those zones moved "
        "to ECDSA. But cumulative adoption cannot fall, so Rogers’ thresholds "
        "would file RSA/SHA-1 as “late majority” in 2011 and “innovator” in 2026 — "
        "the same mechanism walking backwards through the categories.", size=15)

# =========================================================================== #
# 5. Our definition
# =========================================================================== #
s, top = new("Three stages, each an operation on the data",
             "Named after the arithmetic, so no inference can hide inside the name.",
             eyebrow="the definition we use")
cards(s, top, [
    ("STAGE 1 — FIRST SEEN",
     "Present on ≥ 1 zone,\nany RIR",
     "somebody did it once ·\nalways quoted with n"),
    ("STAGE 2 — PARTIAL USAGE",
     "P(value | signed delegations)\n≥ 1%   AND  ≥ 10 zones",
     "in real use, not the norm"),
    ("STAGE 3 — COMMON USAGE",
     "P(value | signed delegations)\n≥ 10%  AND  ≥ 10 zones",
     "a normal operator choice"),
], height=2.1)
textbox(s, M, top + Inches(2.5), W - 2 * M, Inches(1.5),
        "The intervals between them are the measurements:\n"
        "onset (publication → first seen) · establishment (first seen → partial) "
        "· ascent (partial → common)", size=16)
textbox(s, M, top + Inches(3.5), W - 2 * M, Inches(0.9),
        "Thresholds were swept, not picked. The qualifying count is flat across "
        "0.5–3% and 4–25%, so any value mid-plateau gives identical results. "
        "The ≥10-zone guard exists because a population grows: on the reverse "
        "panel one zone was 3.1% of it in 2011 and is 0.016% now — a 201× change "
        "in what the same threshold demands.", size=14, color=MUTED)

# =========================================================================== #
# 5b. The stages on real mechanisms — examples and justification
# =========================================================================== #
s, top = new("The same three stages, three different endings",
             "All three first appeared. Two reached real use, one never did, and one of "
             "those two has since fallen back.",
             eyebrow="worked examples")
after = table(s, M, top, W - 2 * M, [
    ["From RFC 8080", "Onset", "1 · first seen", "2 · in use", "3 · widely used", "Peak", "Now"],
    ["Ed25519 — forward", "1.9 y", "2019-01", "2020-04", "never", "2.04%", "0.003%"],
    ["Ed25519 — reverse", "5.5 y", "2022-08", "never", "never", "0.39%", "0.37%"],
    ["Ed448 — forward", "3.2 y", "2020-05", "never", "never", "0.00%", "0.00%"],
], [0.24, 0.10, 0.15, 0.13, 0.15, 0.11, 0.11], row_h=0.52)

textbox(s, M, after + Inches(0.26), Inches(5.7), Inches(0.34),
        "Why the stages are separate", size=15, bold=True, color=TEAL)
textbox(s, M, after + Inches(0.64), Inches(5.7), Inches(1.9),
        "Ed25519 and Ed448 come from the SAME RFC, published the same day, in the "
        "same family — one date for “RFC 8080” would have averaged them into a "
        "single number." + NL + NL +
        "Ed25519 appeared in the forward TLDs after 1.9 years and in reverse "
        "delegations only after 5.5, and did something different in each. "
        "Ed448 has never left first sighting anywhere." + NL + NL +
        "Each row is one corpus throughout. This is why the unit is the "
        "observable change and not the document, and why a first-sighting date "
        "has to name the corpora it searched.",
        size=13.5)

textbox(s, M + Inches(6.3), after + Inches(0.26), Inches(5.4), Inches(0.34),
        "Why these thresholds", size=15, bold=True, color=TEAL)
textbox(s, M + Inches(6.3), after + Inches(0.64), Inches(5.4), Inches(1.9),
        "Swept, not picked. The number of qualifying algorithms is flat across "
        "0.5–3% and again across 4–25%, so any value mid-plateau gives identical "
        "results. 4% or 30% would sit on a cliff." + NL + NL +
        "The ≥10-zone guard: on the reverse panel, RSA/SHA-256 and RSASHA1-NSEC3 "
        "both “reach 1%” in 2011-05 on a SINGLE zone — 1% of a 32-zone "
        "population." + NL + NL +
        "A ≥25 guard over-corrects, pushing RSASHA1-NSEC3 from 2.0y to 7.8y "
        "purely because it had few contemporaries.",
        size=13.5)

# =========================================================================== #
# 5c. Forward vs reverse -- where two independent corpora agree
# =========================================================================== #
s, top = new("Two corpora that share nothing, compared at one date",
             "2023-12, the last month both cover.",
             eyebrow="cross-reference")
picture(s, "forward_vs_reverse", top, max_h=Inches(3.8))
textbox(s, M, H - Inches(1.6), W - 2 * M, Inches(1.25),
        "Of 20 observables both corpora can answer: 1 agrees, 6 disagree by more "
        "than 5 points, 3 are present on one side only, and 10 read under 0.5% on "
        "both — absent, not corroborating. The one agreement is worth having: "
        "SHA-256 DS digest reads 98.34% and 98.24%, a tenth of a point apart, from "
        "corpora sharing no infrastructure, operators or method. Every "
        "disagreement points the same way: forward zones have modernised and "
        "reverse delegations have not.", size=13.5)

# =========================================================================== #
# 5d. Why they are never pooled
# =========================================================================== #
s, top = new("Why the two are never pooled",
             "The forward sources are 99.6% of the population, and they stop.",
             eyebrow="the denominator")
picture(s, "composition_break", top, max_h=Inches(3.5))
textbox(s, M, H - Inches(1.6), W - 2 * M, Inches(1.2),
        "A pooled share is a forward number until 2023-12 and a reverse number "
        "after it, with a 99.6% cliff between. Peak and latest values routinely "
        "land on opposite sides, so 11 of 15 peak-to-latest drops in the pooled "
        "run are not comparable and none is quoted here. Every figure on this "
        "deck is one corpus at a time.", size=14)

# =========================================================================== #
# 5e. The limitation that matters most
# =========================================================================== #
s, top = new("A share of names is not a count of decisions",
             "In the forward corpus, single operators move six figures of names "
             "in one month.",
             eyebrow="the biggest caveat")
after = table(s, M, top, Inches(5.6), [
    ["Ed25519, forward TLDs", "names", "share"],
    ["2020-01", "3", "0.000%"],
    ["2020-06", "19,863", "2.039%"],
    ["2021-06", "42", "0.004%"],
    ["2023-01", "11,921", "0.557%"],
    ["2023-12", "72", "0.003%"],
], [0.40, 0.30, 0.30], row_h=0.40)
textbox(s, M + Inches(6.1), top, Inches(5.6), Inches(3.4),
        "Tens of thousands of names appearing and vanishing within months, "
        "twice, is not diffusion. It is one operator moving a portfolio — .se "
        "accounts for nearly all of it, peaking at 19,455 names." + NL + NL +
        "The largest jump anywhere is ECDSA P-256: +118,961 names in 2019-02, "
        "76% of its total that month. So “ECDSA reached common usage in 2019-02” "
        "records one bulk migration." + NL + NL +
        "Forward shares therefore largely measure registrar and provider "
        "defaults, not independent decisions.", size=13.5)
textbox(s, M, top + Inches(3.6), W - 2 * M, Inches(0.9),
        "This strengthens the case against the diffusion literature rather than "
        "weakening it: that literature assumes a population of independent "
        "adopters, and this one has a handful of actors who can move six figures "
        "of names at once.", size=14, color=MUTED)

# =========================================================================== #
# 6. The pipeline — what it does
# =========================================================================== #
s, top = new("The pipeline", "Five stages, each resumable, each skippable.",
             eyebrow="what it does")
after = table(s, M, top, W - 2 * M, [
    ["Stage", "What happens", "Output"],
    ["ripe", "Fetches RIPE’s reverse-delegation archive (the only stage that "
             "touches the network)", "the reverse corpus"],
    ["index", "Walks every cache root on the NAS, merges them into one view",
     "inventory.json"],
    ["extract", "One pass per source-day → tidy counts, one checkpoint each",
     "timeline_monthly.csv"],
    ["analyse", "Bottom-up over observable changes, top-down over categories",
     "bottom_up.json · top_down.json"],
    ["report", "Digest for humans, compact bundle for the next conversation",
     "summary.md · analysis_bundle.json"],
], [0.12, 0.55, 0.33], row_h=0.56)
textbox(s, M, after + Inches(0.2), W - 2 * M, Inches(1.2),
        "Both analysis directions are configuration, not code — data/analysis_config.json "
        "holds the observable changes, the implementation groups, the categories and "
        "the thresholds. Adding a mechanism to measure means adding a row.", size=15)

# =========================================================================== #
# 7. The pipeline — how it works
# =========================================================================== #
s, top = new("Three design decisions that protect the numbers",
             eyebrow="how it works")
cards(s, top, [
    ("MULTI-DRIVE, MERGED",
     "The cache spans a main path and a spill path. A day’s files can sit on "
     "both.",
     "Reading one path alone reports a PARTIAL day as a complete one — "
     "indistinguishable from a month when fewer zones signed."),
    ("EVERY SHARE CARRIES ITS DENOMINATOR",
     "Each dimension emits its own _total row alongside the values.",
     "The same fact reads 68%, 64% or 0.6% depending only on what you divide by. "
     "P(x | population) is recomputable, never guessed."),
    ("SHARDED ACROSS MACHINES",
     "Same command, different --shard. Balanced on bytes, no coordinator.",
     "3 shards over the corpus: 0.0% spread, and identical results to a "
     "single-machine run across 30 changes × 7 fields."),
], height=3.3, size=13.5)
textbox(s, M, top + Inches(3.6), W - 2 * M, Inches(0.8),
        "A partial corpus produces entirely plausible numbers, so each shard "
        "reports when it finishes and the analysis refuses to stay silent if the "
        "set is incomplete.", size=14, color=MUTED)

# =========================================================================== #
# 7b. What the full corpus changed
# =========================================================================== #
s, top = new("What the full corpus changed",
             "Adding 8.24 TB of forward data moved 14 dates and refuted one finding.",
             eyebrow="the full run")
cards(s, top, [
    ("MEASURABLE OBSERVABLES",
     "14  →  24",
     "Ten became measurable for the first time. The forward corpus carries NSEC3 "
     "parameters, DNSKEY flags and CDS, which delegation data cannot."),
    ("FIRST-SEEN DATES MOVED EARLIER",
     "14 of them",
     "Existence is a minimum over all evidence, so more corpus can only move a "
     "date earlier. Largest: SHA-384 digest −4.1 y, Ed25519 −3.7 y."),
    ("A FINDING WITHDRAWN",
     "no longer true",
     "“Each new algorithm takes longer to appear” held in reverse data and fails "
     "on the full corpus."),
], height=1.8, size=16)
after = table(s, M, top + Inches(1.98), W - 2 * M, [
    ["New cryptographic primitive", "Published", "Reverse only", "Full corpus"],
    ["ECC-GOST", "2010-07", "2.5 y", "2.4 y"],
    ["ECDSA P-256", "2012-04", "3.7 y", "3.6 y"],
    ["ECDSA P-384", "2012-04", "4.0 y", "3.9 y"],
    ["Ed25519", "2017-02", "5.6 y", "1.9 y  ← the newest is the fastest"],
    ["Ed448", "2017-02", "5.8 y", "3.2 y"],
], [0.30, 0.16, 0.18, 0.36], row_h=0.34)
textbox(s, M, after + Inches(0.12), W - 2 * M, Inches(0.6),
        "Reverse only: 3.7 → 4.0 → 5.6 → 5.8, monotonic. Full corpus: "
        "2.4 → 3.6 → 3.9 → 1.9 → 3.2, not monotonic. The pattern was a property of "
        "one corpus: reverse-DNS operators adopt late, so every newer algorithm "
        "looked slower. .se had Ed25519 in 2019-01; the RIRs not until 2022-08.",
        size=12.5, color=MUTED)

# =========================================================================== #
# 8. Insight — onset does not predict spread
# =========================================================================== #
s, top = new("Being early tells you nothing about ending up used",
             eyebrow="insight 1")
picture(s, "onset_vs_spread", top, max_h=Inches(4.2))
textbox(s, M, H - Inches(1.45), W - 2 * M, Inches(1.1),
        "Peak share, not today's — “today” is 2023-12 for the forward-only "
        "observables and 2026-08 for the DS ones, so one axis cannot carry both. "
        "At n=11 the correlation is not distinguishable from zero under either "
        "endpoint (−0.14 on peak, +0.33 on latest). This is the evidence FOR "
        "splitting the definition, not an assumption behind it.", size=14)

# =========================================================================== #
# 9. Insight — implementation cost predicts onset
# =========================================================================== #
s, top = new("What does predict onset: how many parties must ship code",
             eyebrow="insight 2")
after = table(s, M, top, W - 2 * M, [
    ["What the change requires", "Measured on", "Onset", "n"],
    ["A new codepoint — the cryptography is already linked in",
     "algorithms 8, 10", "0.4 – 0.8 y", "2"],
    ["A new DS digest — a hash, parent side only", "digests 3, 4", "0.8 – 1.2 y", "2"],
    ["Same cryptography, new signalling", "algorithm 7", "1.3 y", "1"],
    ["A new record type on existing infrastructure", "CDS delete", "1.4 y", "1"],
    ["A NEW PRIMITIVE — signer AND validator must both ship",
     "algorithms 12–16", "1.9 – 3.9 y", "5"],
], [0.50, 0.20, 0.20, 0.06], row_h=0.52)
textbox(s, M, after + Inches(0.22), W - 2 * M, Inches(1.7),
        "Five groups, strictly ascending, on the full corpus. The claim to make is "
        "about the ORDERING: the middle gaps are 0.1 y with n=1, which is not a "
        "boundary. Only the last gap is substantial — 1.4 y to 1.9 y, where a "
        "change starts needing both ends." + NL + NL +
        "Not difficulty: Ed25519 is no harder to implement than RSA/SHA-512. "
        "RFC 5218 calls it “incrementally deployable”, which is the citation for "
        "the mechanism.", size=14.5)

# =========================================================================== #
# 10. Insight — displacement
# =========================================================================== #
s, top = new("Most of what happens is retreat, not arrival",
             "Reverse panel, the one series with no composition break: four of "
             "six mechanisms sit below their own peak.",
             eyebrow="insight 3")
picture(s, "displacement", top, max_h=Inches(3.7))
textbox(s, M, H - Inches(1.5), W - 2 * M, Inches(1.1),
        "Every framework in the literature describes something arriving. None "
        "describes something being pushed out. In a fixed population that is half "
        "of what happens — and in our data it is the bigger half. This looks like "
        "the contribution.", size=15, bold=False)

# =========================================================================== #
# 11. Insight — deprecation
# =========================================================================== #
s, top = new("Deprecation runs on an inverted clock",
             "The question is not when it appeared but how much is still there.",
             eyebrow="insight 4")
cards(s, top, [
    ("RFC 9905 — SHA-1 SIGNING",
     "7.30%\n15,888 records",
     "still published after the RFC forbade new use"),
    ("RFC 9905 — SHA-1 DS DIGEST",
     "18.26%\n19,652 records",
     "a different fix: replace the DS at the parent"),
    ("RFC 9906 — ECC-GOST",
     "0\nzero records",
     "documented an ending rather than causing one"),
], height=2.2, size=20)
textbox(s, M, top + Inches(2.7), W - 2 * M, Inches(1.4),
        "The two halves of RFC 9905 need different remedies — replace the DS at "
        "the parent, versus reissue the child’s keys — so they are reported as two "
        "numbers, never one “SHA-1 exposure”." + NL + NL +
        "Both RFCs published 2025-11 and the forward corpus ends 2023-12, so the "
        "entire post-publication window is reverse data — a property of the "
        "corpus, not a choice.", size=15)

# =========================================================================== #
# 12. Limits
# =========================================================================== #
s, top = new("What this cannot see", "Stated so no one reads a blank as a zero.",
             eyebrow="limits")
table(s, M, top, W - 2 * M, [
    ["Limit", "Consequence"],
    ["Forward and reverse cover different periods",
     "Forward ends 2023-12, the RIRs run to 2026-08. Any share spanning 2024-01 "
     "compares two populations: 11 of 15 peak-to-now drops are not comparable, "
     "and none of them is quoted."],
    ["Only TLSA is still unreachable",
     "The forward corpus closed the rest — 2 observables have no denominator "
     "anywhere, down from 9. Untestable, not negative."],
    ["We read published zone data, never live queries",
     "availability / verifiability / validity (Osterweil’s three) are outside "
     "the instrument entirely."],
    ["Nothing in a zone file records intent",
     "We never claim an RFC CAUSED a change. An operator whose provider changed "
     "a default is indistinguishable from one who read the RFC."],
    ["Our checklist covers 29 of 105 current DNSSEC RFCs",
     "13 of the uncovered ones are measurable from zone data — ZONEMD, "
     "compact denial of existence, multi-signer, OPENPGPKEY, SSHFP."],
], [0.34, 0.66], row_h=0.78)

# =========================================================================== #
# 13. What needs to be done
# =========================================================================== #
s, top = new("What needs to be done", eyebrow="next")
table(s, M, top, W - 2 * M, [
    ["", "Work", "Why it matters"],
    ["1", "Analyse forward and reverse separately, never pooled",
     "They cover different periods, so a pooled share crosses a break at "
     "2024-01. This is the one thing blocking the prevalence numbers; it needs "
     "no rescan, only the analyse stage."],
    ["2", "Extend the forward mirror past 2023-12",
     "RFC 9905 and 9906 published 2025-11, so the whole deprecation question is "
     "answerable only on reverse data until the forward mirror catches up."],
    ["3", "Decide: share measure, cumulative measure, or both",
     "A cumulative measure would let us inherit Rogers’ thresholds, but it "
     "cannot see displacement."],
    ["4", "Settle whether displacement is genuinely unclaimed",
     "If it is, it is a stronger contribution than the timeline. Needs someone "
     "who knows the literature better than a search does."],
    ["5", "Add the 13 measurable RFCs we do not screen",
     "ZONEMD, compact denial of existence, multi-signer. Mechanical work that "
     "would extend the story rather than just widen it."],
], [0.04, 0.36, 0.60], row_h=0.78)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {OUT}")

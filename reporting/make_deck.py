"""Assemble the DNSSEC RFC-adoption deck."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

CHARTS = Path(sys.argv[1])
OUT = Path(sys.argv[2])

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
CRITICAL = RGBColor(0xD0, 0x3B, 0x3B)
RULE = RGBColor(0xE1, 0xE0, 0xD9)
FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SURFACE
    return s


def textbox(s, x, y, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
            spacing=1.0, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT
    return tb


def rule(s, y, x=M, w=None, color=RULE, h=Pt(1.2)):
    w = w or (W - 2 * M)
    shp = s.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def header(s, kicker, title, insight=None):
    """Kicker + title + the takeaway. The insight line carries the point."""
    textbox(s, M, Inches(0.52), W - 2 * M, Inches(0.3), kicker.upper(), 12,
            MUTED, bold=True)
    textbox(s, M, Inches(0.86), W - 2 * M, Inches(0.7), title, 30, INK, bold=True)
    if insight:
        textbox(s, M, Inches(1.62), W - 2 * M, Inches(0.6), insight, 15, INK2,
                spacing=1.25)


def footnote(s, text):
    textbox(s, M, H - Inches(0.62), W - 2 * M, Inches(0.35), text, 10.5, MUTED)


def chart(s, name, top=Inches(2.35), height=Inches(4.35)):
    img = CHARTS / name
    from PIL import Image
    with Image.open(img) as im:
        ar = im.width / im.height
    h = height
    w = Emu(int(h * ar))
    if w > W - 2 * M:
        w = W - 2 * M
        h = Emu(int(w / ar))
    s.shapes.add_picture(str(img), Emu(int((W - w) / 2)), top, width=w, height=h)


def stat_row(s, items, y=Inches(2.6), color=INK):
    """Hero figures. Proportional figures, sans face, value + label."""
    n = len(items)
    gap = Inches(0.4)
    cw = Emu(int((W - 2 * M - gap * (n - 1)) / n))
    for i, (value, label) in enumerate(items):
        x = Emu(int(M + i * (cw + gap)))
        # Long values (a year range) do not fit a column at hero size; step down
        # rather than let the text overflow its box.
        size = 40 if len(value) <= 6 else 30
        textbox(s, x, y, cw, Inches(0.78), value, size, color, bold=True)
        textbox(s, x, y + Inches(0.86), cw, Inches(0.9), label, 13, INK2, spacing=1.2)


# ----------------------------------------------------------------- 1. title --
s = slide()
textbox(s, M, Inches(2.25), W - 2 * M, Inches(0.4), "OPENINTEL RFC-ADOPTION ANALYSIS",
        13, MUTED, bold=True)
textbox(s, M, Inches(2.72), W - 2 * M, Inches(1.4),
        "How DNSSEC RFC adoption\nactually changed, 2018–2026", 42, INK, bold=True,
        spacing=1.08)
rule(s, Inches(4.5), w=Inches(2.2), color=BLUE, h=Pt(3))
textbox(s, M, Inches(4.85), Inches(9.5), Inches(1.0),
        "2.76 billion DNS records  ·  3,127 measurement days  ·  3 zones  ·  8 RFCs\n"
        "Evidence-linked: every count traces to a matched indicator and a publication-date check.",
        14, INK2, spacing=1.35)

# ------------------------------------------------------- 2. what was measured --
s = slide()
header(s, "Scope", "What was measured",
       "Every DNS record was tested against all 8 RFC checklists, not one chosen in advance.")
stat_row(s, [
    ("2.76B", "DNSSEC records evaluated\nafter the record-type prefilter"),
    ("3,127", "measurement days\n.gov, .nu, .se"),
    ("8", "DNSSEC RFCs\nscored independently"),
    ("2018–2026", "observation window\n(2022 and 2025 absent)"),
])
rule(s, Inches(4.9))
textbox(s, M, Inches(5.2), W - 2 * M, Inches(1.4),
        "Each record is checked against every RFC's observable signature, then against that RFC's\n"
        "publication date — an observation cannot evidence a standard that did not yet exist.\n"
        "Counts below are exact aggregates; rankings are derived from a sampled set of worked examples.",
        14, INK2, spacing=1.4)
footnote(s, "Source: OpenINTEL forward-DNS, basis=zonefile. Analysis: openintel_rfc pipeline.")

# ---------------------------------------------------- 3. read this first ------
s = slide()
header(s, "Read this first", "The three zones are not equally sampled",
       "Totals across all zones are dominated by whichever large zone happened to be measured. "
       "Every later chart is normalised within a zone.")
chart(s, "panel_balance.png", top=Inches(2.5), height=Inches(3.5))
textbox(s, M, Inches(6.2), W - 2 * M, Inches(0.6),
        ".gov is 70% of the days but 2.6% of the records.  .se is 1.2% of the days but 20.6% of the records.",
        13, CRITICAL, bold=True)
footnote(s, "2020 and 2026 contain .gov only — an all-zone time series would show a collapse that is purely sampling.")

# --------------------------------------------------------- 4. THE headline ----
s = slide()
header(s, "The headline", "ECDSA replaced RSA as the default",
       "Monotonic over eight years, and real growth rather than composition: 1,035 to 30,284 records per day in .gov, a 29x rise against a 1.8x baseline.")
chart(s, "ecdsa_migration.png", top=Inches(2.3), height=Inches(4.3))
footnote(s, "RFC 6605 (algorithms 13/14), share of each zone's own DNSSEC records. 2022 and 2025 not measured.")

# ------------------------------------------------- 5. cross-zone, common year --
s = slide()
header(s, "Cross-zone", "The same year, three zones",
       "In 2021 the ccTLDs led .gov on ECDSA — .gov overtook them only later. "
       "Comparing each zone's latest year instead would invert this.")
chart(s, "landscape.png", top=Inches(2.4), height=Inches(4.2))
footnote(s, "2021 is the only year all three zones were measured. .se rests on 11 days — directional, not precise.")

# --------------------------------------------------------- 6. EdDSA ----------
s = slide()
header(s, "The non-adoption", "EdDSA was standardised in 2017. It never arrived.",
       "Nine years on, adoption is indistinguishable from zero in two of three zones.")
stat_row(s, [
    ("41", ".gov records using EdDSA\nout of 72 million"),
    ("0.09%", "of .nu DNSSEC records"),
    ("0.84%", "of .se DNSSEC records\n— 72% of all EdDSA seen"),
], color=CRITICAL)
rule(s, Inches(4.9))
textbox(s, M, Inches(5.2), W - 2 * M, Inches(1.3),
        "Registries consolidated on ECDSA instead. A newer, smaller, faster algorithm with IETF backing\n"
        "is not sufficient for adoption — ECDSA arrived first and was good enough.\n"
        ".se is the only zone with a measurable EdDSA presence.",
        14, INK2, spacing=1.4)
footnote(s, "RFC 8080 (algorithms 15/16), published February 2017. Counts across the full 2.76 billion records.")

# ------------------------------------------- 7. .gov delegation automation ----
s = slide()
header(s, "Operational outlier", ".gov automates delegation; the ccTLDs do not",
       "CDS/CDNSKEY publication grew 11× and the delete signal appeared in 2020 and kept climbing.")
chart(s, "gov_automation.png", top=Inches(2.4), height=Inches(4.05))
footnote(s, ".gov holds 56% of all delete signals and 64% of all CDS/CDNSKEY records while being 2.6% of the data.")

# ------------------------------------------------ 8. what shares hide --------
s = slide()
header(s, "What shares hide", "NSEC3 did not decline. It stood still.",
       "Its share of .gov records fell 10.9% to 5.8% — but per measurement day it held flat "
       "at ~2,900 for eight years. The share moved because the denominator grew 79%.")
chart(s, "growth_vs_baseline.png", top=Inches(2.45), height=Inches(3.95))
footnote(s, "Per-day rates remove the shared denominator. Only growth above the zone's own 1.8x baseline is adoption.")

# --------------------------------------------------- 9. what this is not ------
s = slide()
header(s, "Limits", "What this does not show",
       "Stated plainly, because each one changes how a number should be read.")
items = [
    ("Record-level, not zone-level.",
     "“62% of .gov records use ECDSA” is not “62% of .gov zones”. Large signed zones contribute more records."),
    ("First-seen is a coverage floor.",
     "Five RFCs report first-seen 2018-01-01 — the first day of data. They were deployed before the window."),
    ("Rankings are sampled; counts are not.",
     "Observation counts are exact aggregates. Scores derive from 29 worked examples, ~1 per 95M rows."),
    ("Presence is not conformance.",
     "A matching record shows the mechanism is deployed, not that the operator read the RFC."),
]
y = Inches(2.5)
for head, body in items:
    textbox(s, M, y, Inches(4.4), Inches(0.5), head, 15, INK, bold=True)
    textbox(s, M + Inches(4.7), y, W - 2 * M - Inches(4.7), Inches(0.8), body, 14, INK2,
            spacing=1.3)
    y += Inches(1.02)
footnote(s, "The pipeline identifies ranked RFC candidates consistent with observable signals; it does not prove adoption.")

# --------------------------------------------------------- 10. next ----------
s = slide()
header(s, "Next", "What would make this publishable",
       "Three gaps, in the order they change the conclusions.")
items = [
    ("1", "Aggregate to zone level",
     "Turns “share of records” into “share of zones” — the number a reader assumes they are already seeing."),
    ("2", "Fill 2022 and 2025, and sample .se properly",
     ".se is the largest zone and the worst covered at 39 days. The ECDSA curve needs the missing years."),
    ("3", "Exact scoring in SQL",
     "Removes the sampled-ranking caveat so scores are as exact as the counts already are."),
]
y = Inches(2.55)
for num, head, body in items:
    textbox(s, M, y - Inches(0.06), Inches(0.6), Inches(0.6), num, 26, BLUE, bold=True)
    textbox(s, M + Inches(0.72), y, Inches(4.6), Inches(0.5), head, 16, INK, bold=True)
    textbox(s, M + Inches(5.5), y, W - 2 * M - Inches(5.5), Inches(0.9), body, 14, INK2,
            spacing=1.3)
    y += Inches(1.25)
rule(s, Inches(6.35))
textbox(s, M, Inches(6.6), W - 2 * M, Inches(0.5),
        "Every figure in this deck is reproducible: openintel-rfc merge --checkpoint-dir out/final/checkpoints",
        12, MUTED)

prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

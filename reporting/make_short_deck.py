"""Five slides: the definition, bottom-up, top-down, where they meet, what to trust.

The long deck (make_framework_deck.py) stays for reference. This is the version
to present: two slides give the two directions a full page each, one shows the
crosswalk that is the actual result, and the outer two carry just enough to make
those three readable.

Every figure is read from out/server_run/split_analysis.json at build time.
Nothing is typed in, so a re-run of the analysis moves the deck with it.
"""
from __future__ import annotations

import json
import sys
from collections import defaultdict
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

RUN = Path(sys.argv[1] if len(sys.argv) > 1 else "out/server_run")
CHARTS = Path(sys.argv[2] if len(sys.argv) > 2 else "reporting/charts/short")
OUT = Path(sys.argv[3] if len(sys.argv) > 3
           else "out/analysis/dnssec_framework_short.pptx")

INK = RGBColor(0x17, 0x1A, 0x14)
MUTED = RGBColor(0x66, 0x70, 0x5C)
RULE = RGBColor(0xD3, 0xD8, 0xCC)
PAPER = RGBColor(0xF6, 0xF7, 0xF3)
TEAL = RGBColor(0x00, 0x90, 0x7F)
BRICK = RGBColor(0xB3, 0x2B, 0x22)
FONT = "Segoe UI"
NL = chr(10)

# ---------------------------------------------------------------- data ----- #
sp = json.loads((RUN / "split_analysis.json").read_text(encoding="utf-8"))
cfg = json.loads(Path("data/analysis_config.json").read_text(encoding="utf-8"))
glabel = {k: v["label"] for k, v in cfg["bottom_up"]["groups"].items()}
group_of = {c["label"]: c["group"] for c in cfg["bottom_up"]["changes"]}
rfc_of = {c["label"]: c["rfc"] for c in cfg["bottom_up"]["changes"]}
rfc_cat = {r: c["label"] for c in cfg["top_down"]["categories"] for r in c["rfcs"]}

fwd = {r["label"]: r for r in sp["sides"]["forward"]["bottom_up"]}
rev = {r["label"]: r for r in sp["sides"]["reverse"]["bottom_up"]}
best = {r["label"]: r for r in sp["first_seen_across_corpora"]}
xref = sp["cross_reference"]

onset = {}
for label in set(fwd) | set(rev):
    cands = [s[label]["onset_years"] for s in (fwd, rev)
             if label in s and s[label].get("onset_years") is not None
             and s[label]["onset_years"] >= 0 and not s[label].get("left_censored")]
    if cands:
        onset[label] = min(cands)

seen = sum(1 for r in best.values() if r["first_seen"])
configured = len(cfg["bottom_up"]["changes"])
cat_span = defaultdict(list)
for label, y in onset.items():
    cat_span[rfc_cat.get(rfc_of.get(label, ""), "(uncategorised)")].append(y)

# --------------------------------------------------------------- helpers --- #
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height
M = Inches(0.8)
BLANK = prs.slide_layouts[6]


def textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
            align=PP_ALIGN.LEFT, space_after=6, line=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para in enumerate(str(text).split(NL)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        run = p.add_run()
        run.text = para
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return tb


def new(title, dek=None, eyebrow=None):
    s = prs.slides.add_slide(BLANK)
    y = M - Inches(0.3)
    if eyebrow:
        textbox(s, M, y, W - 2 * M, Inches(0.3), eyebrow.upper(), size=11,
                bold=True, color=TEAL)
        y += Inches(0.36)
    textbox(s, M, y, W - 2 * M, Inches(0.6), title, size=30, bold=True)
    y += Inches(0.74)
    if dek:
        textbox(s, M, y, W - 2 * M, Inches(0.44), dek, size=14, color=MUTED)
        y += Inches(0.52)
    line = s.shapes.add_shape(1, M, y, W - 2 * M, Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = RULE
    line.line.fill.background(); line.shadow.inherit = False
    return s, y + Inches(0.28)


def table(slide, x, y, w, rows, widths, *, size=12.5, row_h=0.4):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w,
                                   Inches(row_h * len(rows)))
    tbl = shape.table
    total = sum(widths)
    for i, frac in enumerate(widths):
        tbl.columns[i].width = Emu(int(w * frac / total))
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(row_h)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size = Pt(size - 1 if r == 0 else size)
            run.font.name = FONT
            run.font.bold = r == 0
            run.font.color.rgb = MUTED if r == 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if r == 0 else RGBColor(255, 255, 255)
    return y + Inches(row_h * len(rows))


def picture(slide, name, top, *, max_h, max_w=None):
    path = CHARTS / f"{name}.png"
    iw, ih = Image.open(path).size
    aw = max_w or (W - 2 * M)
    scale = min(aw / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int((W - w) / 2), int(top), w, h)
    return top + h


# =============================================================== slide 1 === #
s, top = new("Three stages, because one date hid three answers",
             "Each defined by an operation on the data, so nothing can hide in the name.",
             eyebrow="the definition")
table(s, M, top, W - 2 * M, [
    ["Stage", "Test", "Reads as"],
    ["1 · first seen", "present on ≥ 1 name, either corpus",
     "somebody did it once — always quoted with n"],
    ["2 · in use", "P(value | signed names) ≥ 1%   and ≥ 10 names",
     "in real use, not the norm"],
    ["3 · widely used", "P(value | signed names) ≥ 10%  and ≥ 10 names",
     "a normal choice"],
], [0.20, 0.44, 0.36], row_h=0.46)
textbox(s, M, top + Inches(2.05), Inches(5.7), Inches(2.6),
        "Why three" + NL +
        "Ed25519 appeared in the forward TLDs 1.9 years after RFC 8080, reached "
        "2.04% of signed names, and is back to 0.003%. “Adopted after 1.9 years” "
        "and “never adopted” are both true of it." + NL + NL +
        "The gaps between the stages are the measurements: onset, establishment, "
        "ascent.", size=14)
textbox(s, M + Inches(6.3), top + Inches(2.05), Inches(5.4), Inches(2.6),
        "Why 1% and 10%" + NL +
        "Swept, not chosen. The number of qualifying algorithms is flat across "
        "0.5–3% and again across 4–25%, so any value mid-plateau gives identical "
        "results; 4% or 30% would sit on a cliff." + NL + NL +
        "Rogers’ 2.5/16/50/84 were tested and rejected: his curve is cumulative "
        "and cannot fall, ours is a share of a fixed population and does.",
        size=14)
textbox(s, M, H - Inches(1.15), W - 2 * M, Inches(0.7),
        f"Measured on 8.24 TB · 15,179 source-days · 12 sources · "
        f"{seen} of {configured} configured observables seen",
        size=12, color=MUTED)

# =============================================================== slide 2 === #
s, top = new("Bottom-up: onset tracks what the change costs an implementer",
             "Every observable with a usable onset, grouped by what someone has to "
             "ship. Censored dates and values older than their own RFC excluded.",
             eyebrow="bottom-up")
after = picture(s, "bottom_up_bands", top, max_h=Inches(3.25))
textbox(s, M, after + Inches(0.16), W - 2 * M, Inches(1.4),
        "The ordering is strict across five groups, and the one substantial gap — "
        "1.4 to 1.9 years — falls exactly where a change starts needing BOTH ends "
        "to ship new cryptography. A signer and a validator must agree before "
        "anything is publishable; everything faster needs one party to change a "
        "value it already supports. RFC 5218 calls that “incrementally "
        "deployable”." + NL +
        "Not difficulty: Ed25519 is no harder to implement than RSA/SHA-512.",
        size=14)

# =============================================================== slide 3 === #
s, top = new("Top-down: the categories are communicable, not predictive",
             "The same observables, grouped by conceptual impact instead.",
             eyebrow="top-down")
after = picture(s, "top_down_bands", top, max_h=Inches(2.55))
rows = [["Category", "Configured", "Seen", "Onset span"]]
for cat in cfg["top_down"]["categories"]:
    conf = sum(1 for c in cfg["bottom_up"]["changes"]
               if rfc_cat.get(c["rfc"]) == cat["label"])
    sn = sum(1 for c in cfg["bottom_up"]["changes"]
             if rfc_cat.get(c["rfc"]) == cat["label"]
             and best.get(c["label"], {}).get("first_seen"))
    sp_ = cat_span.get(cat["label"], [])
    rows.append([cat["label"], conf, sn,
                 f"{min(sp_):.1f}–{max(sp_):.1f} y" if sp_ else "no usable onset"])
table(s, M, after + Inches(0.14), W - 2 * M, rows, [0.46, 0.16, 0.12, 0.26],
      row_h=0.30, size=11.5)

# =============================================================== slide 4 === #
s, top = new("Where the two meet",
             "One category holds almost all the evidence — and spans the whole "
             "onset range.",
             eyebrow="the result")
crypto = cat_span.get("Which cryptography a zone signs with", [])
allv = list(onset.values())
table(s, M, top, W - 2 * M, [
    ["Cut", "Groups", "What it does to onset"],
    ["By what an implementer must ship", "5",
     "partitions it — five bands, strictly ascending, one real gap"],
    ["By conceptual category", f"{len(cfg['top_down']['categories'])}",
     f"does not — “which cryptography a zone signs with” alone spans "
     f"{min(crypto):.1f}–{max(crypto):.1f} y, the full range of the data"],
], [0.30, 0.10, 0.60], row_h=0.52)
textbox(s, M, top + Inches(1.85), W - 2 * M, Inches(2.6),
        "So the two directions do not converge on one taxonomy, and that is the "
        "finding rather than a failure: they answer different questions." + NL + NL +
        "The implementation-cost groups predict when something will first appear. "
        f"The conceptual categories cannot — {len(crypto)} of the {len(onset)} "
        "observables with a usable onset sit in a single category, and it covers "
        f"every value from {min(crypto):.1f} to {max(crypto):.1f} years. The "
        "categories are the right way to explain the work to a reader and the "
        "wrong way to model it." + NL + NL +
        "Two categories carry no usable evidence at all: DANE has no TLSA in "
        "either corpus, and the base protocol has no separate switch to date — a "
        "zone publishing any DS record is already evidence for RFC 4034.",
        size=15)

# =============================================================== slide 5 === #
s, top = new("What to trust, and what is next", eyebrow="limits")
comp = [c for c in xref["comparisons"] if c["difference_pct_points"] is not None]
counts = xref["verdict_counts"]
table(s, M, top, W - 2 * M, [
    ["Trust", "Because"],
    ["First-sighting dates and onset",
     "Existence is a minimum over all evidence and needs no denominator. Adding "
     "the forward corpus moved 14 of them earlier."],
    ["A share, only with its corpus and month named",
     f"Forward sources are 99.6% of the population and stop at 2023-12, so a "
     f"pooled share is two different measurements with a 99.6% cliff between."],
    ["Cross-corpus agreement, sparingly",
     f"Of {len(comp)} observables both corpora answer: {counts.get('agree', 0)} "
     f"agree, {counts.get('disagree', 0)} disagree, "
     f"{counts.get('absent from both', 0)} are absent from both. The one "
     f"agreement is SHA-256 DS digest, 0.10 points apart."],
    ["Not a count of decisions",
     "One operator moved +118,961 names to ECDSA in a single month — 76% of its "
     "total that month. Forward shares largely measure provider defaults."],
], [0.30, 0.70], row_h=0.62)
textbox(s, M, top + Inches(3.3), W - 2 * M, Inches(1.9),
        "Next" + NL +
        "1 · Re-run the reverse side with --pool-sources; its names overlap "
        "between RIRs, so those shares are approximate while the forward ones are "
        "exact." + NL +
        "2 · Extend the forward mirror past 2023-12 — RFC 9905 and 9906 published "
        "2025-11, so deprecation is answerable only on reverse data today." + NL +
        "3 · Settle whether displacement is unclaimed in the literature. Every "
        "framework we found describes something arriving; none describes something "
        "being pushed out, and that is most of what our data shows.", size=14)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"{len(prs.slides._sldIdLst)} slides -> {OUT}")

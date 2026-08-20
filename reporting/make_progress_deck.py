"""Progress-update deck: what changed since the last review, and what it found.

Reuses the visual system of `make_deck.py` -- same palette, same helpers, same
rhythm -- so this reads as the next instalment rather than a different document.

Every figure is one already rendered and checked into reporting/charts/, and every
number on a slide is one the cross-reference notebook computes. Nothing here is
retyped from memory.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

CHARTS = Path(sys.argv[1] if len(sys.argv) > 1 else "reporting/charts")
OUT = Path(sys.argv[2] if len(sys.argv) > 2
           else "out/analysis/dnssec_progress_update.pptx")

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
GREEN = RGBColor(0x1B, 0xAF, 0x7A)
CRITICAL = RGBColor(0xD0, 0x3B, 0x3B)
RULE = RGBColor(0xE1, 0xE0, 0xD9)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height
M = Inches(0.85)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = SURFACE
    return s


def textbox(s, x, y, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
            spacing=1.0, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def rule(s, y, x=M, w=None, color=RULE, h=Pt(1.2)):
    w = w if w is not None else W - 2 * M
    box = s.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def header(s, kicker, title, insight=None):
    textbox(s, M, Inches(0.52), W - 2 * M, Inches(0.3), kicker.upper(), 12,
            MUTED, bold=True)
    textbox(s, M, Inches(0.86), W - 2 * M, Inches(0.7), title, 30, INK, bold=True)
    if insight:
        textbox(s, M, Inches(1.62), W - 2 * M, Inches(0.62), insight, 15, INK2,
                spacing=1.25)


def footnote(s, text):
    rule(s, H - Inches(0.86))
    textbox(s, M, H - Inches(0.70), W - 2 * M, Inches(0.5), text, 11, MUTED,
            spacing=1.2)


def chart(s, name, top=Inches(2.40), height=Inches(4.05)):
    img = CHARTS / name
    with Image.open(img) as im:
        ar = im.width / im.height
    h, w = height, Emu(int(height * ar))
    if w > W - 2 * M:
        w = W - 2 * M
        h = Emu(int(w / ar))
    s.shapes.add_picture(str(img), Emu(int((W - w) / 2)), top, width=w, height=h)


def stat_row(s, items, y=Inches(2.6), color=INK):
    n = len(items)
    gap = Inches(0.4)
    cw = Emu(int((W - 2 * M - gap * (n - 1)) / n))
    for i, (value, label) in enumerate(items):
        x = Emu(int(M + i * (cw + gap)))
        size = 40 if len(value) <= 6 else 28
        textbox(s, x, y, cw, Inches(0.78), value, size, color, bold=True)
        textbox(s, x, y + Inches(0.86), cw, Inches(1.0), label, 13, INK2, spacing=1.2)


def bullets(s, items, y=Inches(2.35), step=Inches(0.92), label_w=Inches(4.5)):
    for i, (head, body) in enumerate(items):
        top = Emu(int(y + i * step))
        textbox(s, M, top, label_w, Inches(0.4), head, 15, INK, bold=True)
        textbox(s, Emu(int(M + label_w + Inches(0.3))), top,
                W - 2 * M - label_w - Inches(0.3), Inches(0.8), body, 13.5, INK2,
                spacing=1.22)


# =============================================================== 1. title ====
s = slide()
textbox(s, M, Inches(2.55), W - 2 * M, Inches(0.34), "PROGRESS UPDATE", 13,
        MUTED, bold=True)
textbox(s, M, Inches(2.95), W - 2 * M, Inches(1.5),
        "DNSSEC RFC adoption:\na second corpus, and what it corroborates", 38,
        INK, bold=True, spacing=1.1)
rule(s, Inches(4.62), w=Inches(2.6), color=BLUE, h=Pt(3.5))
textbox(s, M, Inches(5.05), W - 2 * M, Inches(1.3),
        "30 DNSSEC RFCs screened  ·  two independent corpora  ·  2009-2026\n"
        "Every figure traces to a committed checkpoint and an executed notebook.\n"
        "Nothing in this deck is asserted without the cell that computes it.",
        14, INK2, spacing=1.35)

# ============================================================= 2. where we are
s = slide()
header(s, "since the last review", "What changed",
       "Two things: the screening got much wider, and the analysis got a second, "
       "independent corpus to check itself against.")
stat_row(s, [
    ("8 → 30", "DNSSEC RFCs screened.\nDates and status pulled from the RFC Editor\nindex and IANA registries at build time"),
    ("1 → 2", "independent corpora.\nOpenINTEL forward DNS, and RIPE's\nreverse-delegation archive"),
    ("2018 → 2009", "measurement window.\nNine more years, which turns\ncensored bounds into measurements"),
    ("790", "tests passing.\nIncluding cross-engine agreement,\nwhich caught two real bugs"),
], y=Inches(2.55))
rule(s, Inches(4.95))
textbox(s, M, Inches(5.20), W - 2 * M, Inches(1.3),
        "The pipeline also learned to mirror rather than stream, and to pace itself "
        "against the store's rate limit — measured at ~1 request/second with a burst "
        "of five. A nightly script now keeps both corpora current unattended.",
        13.5, INK2, spacing=1.3)
footnote(s, "Checklist 0.2.1 · dictionary 0.2.0 · branch overnight-data")

# ======================================================= 3. the second corpus
s = slide()
header(s, "the new corpus", "RIPE reverse-delegation zones, 2009-2026",
       "Every delegation in in-addr.arpa and ip6.arpa, all five RIRs. It answers two "
       "questions OpenINTEL structurally cannot.")
bullets(s, [
    ("A denominator of zones",
     "Zone files list every delegation, signed or not — so \"share of zones\" is "
     "directly countable. That is the caveat that has been on half our slides."),
    ("Dates before 2018",
     "It predates four of the algorithm RFCs, so their adoption lag is measured "
     "rather than bounded from below by when measurement began."),
    ("A different population",
     "Network operators holding address blocks, not registrants holding domains. "
     "Where it agrees with OpenINTEL, the agreement is evidence."),
], y=Inches(2.60), step=Inches(1.05), label_w=Inches(3.6))
rule(s, Inches(5.95))
textbox(s, M, Inches(6.18), W - 2 * M, Inches(0.6),
        "199 monthly snapshots · 1,875,584 DS records · ~1.3M delegations per day",
        14, BLUE, bold=True)
footnote(s, "Ingested and scanned locally; the archive publishes some days as "
            "zero-byte files, which are treated as gaps rather than as empty days.")

# ========================================================== 4. ECDSA converge
s = slide()
header(s, "corroboration", "ECDSA converges in both corpora",
       "Different infrastructure, different operators, different collection method. "
       "They start 15 points apart in 2018 and land 1.65 points apart in 2026.")
chart(s, "crossref/01_ecdsa_convergence.png", top=Inches(2.42), height=Inches(3.95))
footnote(s, "Forward: share of DNSSEC records across .gov/.nu/.se. Reverse: share of "
            "DS records across all five RIRs. RFC 6605 is matched on algorithm 13/14, "
            "independent of record type, which is what makes the two comparable.")

# ====================================================== 5. denominator lesson
s = slide()
header(s, "a methodology result", "A 10x disagreement that was arithmetic, not behaviour",
       "RFC 4509 reads ~5% forward and ~79% reverse. Normalised to the same "
       "denominator, the two agree to 1.1 points.")
chart(s, "crossref/02_denominator_artefact.png", top=Inches(2.45), height=Inches(4.25))
footnote(s, "The rule we now apply: an indicator scoped by record type is not "
            "comparable across corpora with different record-type composition; one "
            "scoped by algorithm is. Every cross-corpus claim has to say which it is.")

# ============================================================ 6. zone-level
s = slide()
header(s, "the number that answers the caveat", "Zone-level deployment: 0.88-1.01%",
       "After seventeen years. The range is panel choice, not uncertainty in the "
       "data — and whether it has \"crossed 1%\" sits inside it.")
chart(s, "crossref/03_zone_level_share.png", top=Inches(2.40), height=Inches(3.55))
footnote(s, "Two composition breaks are excluded from the headline panel: APNIC leaves "
            "the archive in Jan 2025, and RIPE changes publication format in Oct 2015 "
            "and its delegation count falls 97%. Both would read as adoption if left in.")

# ============================================================ 7. adoption lag
s = slide()
header(s, "measured, not bounded", "Each new signing algorithm took longer to arrive",
       "The forward corpus can only bound these from above. The reverse corpus "
       "starts early enough to measure them.")
stat_row(s, [
    ("0.5 y", "RSA/SHA-2 (RFC 5702)\npublished 2009-10"),
    ("2.5 y", "GOST (RFC 5933)\npublished 2010-07"),
    ("3.7 y", "ECDSA (RFC 6605)\npublished 2012-04"),
    ("5.6 y", "EdDSA (RFC 8080)\npublished 2017-02"),
], y=Inches(2.60), color=BLUE)
rule(s, Inches(4.65))
textbox(s, M, Inches(4.88), W - 2 * M, Inches(1.5),
        "Monotonically increasing by publication date. Two readings, and this data "
        "does not choose between them: either the ecosystem has slowed, or RSA/SHA-2 "
        "was unusually fast because it was a drop-in change to an algorithm operators "
        "already ran, where ECDSA and EdDSA each needed new tooling and new parent "
        "support. The second is more plausible — but it is an interpretation, not a "
        "measurement.",
        14, INK2, spacing=1.3)
footnote(s, "Uncensored: the corpus begins 2009-04, before all four were published. "
            "Five other RFCs remain left-censored and are reported as upper bounds.")

# ================================================================= 8. IPv6
s = slide()
header(s, "the sharpest split in the data", "IPv6 delegations are 9.9x more likely to be signed",
       "8.16% of ip6.arpa delegations carry a DS, against 0.83% of in-addr.arpa. "
       "The gap has been open since 2009 and is widening.")
chart(s, "crossref/06_ipv4_vs_ipv6.png", top=Inches(2.42), height=Inches(3.45))
footnote(s, "Selection, most likely, not causation: an operator who deployed IPv6 "
            "reverse DNS has already done discretionary modern DNS work. It rests on "
            "454 signed delegations — quote the ratio, not the curve's shape.")

# =============================================================== 9. regional
s = slide()
header(s, "who the pooled number is about", "A six-fold spread between regions",
       "LACNIC leads at 5.2% and ARIN trails at 0.84% — but ARIN holds 85% of all "
       "delegations, so the pooled figure is essentially ARIN's.")
chart(s, "crossref/07_by_rir.png", top=Inches(2.50), height=Inches(3.55))
footnote(s, "Any \"reverse DNS is N% signed\" claim is a claim about North American "
            "address space unless it says otherwise — worth stating wherever we quote "
            "the headline. Grey bars carry a composition break in their series.")

# ================================================================ 10. SHA-1
s = slide()
header(s, "operational finding", "\"SHA-1\" is two mechanisms, at different levels",
       "RFC 9905 (Nov 2025) closes both to new deployment. Reporting one number "
       "misstates the exposure by 3x and points at the wrong fix.")
stat_row(s, [
    ("15.3%", "of DS records use a SHA-1 digest (type 1).\nRFC 9905: MUST NOT create new DS.\nFix: replace the DS at the parent."),
    ("4.9%", "point at a SHA-1 signature algorithm (5/7).\nRFC 9905: MUST NOT create new DNSKEY/RRSIG.\nFix: reissue the child's keys."),
    ("0.07%", "still use GOST R 34.11-94.\nRFC 9906 retired it in Nov 2025 —\nafter it had already died out."),
], y=Inches(2.65), color=CRITICAL)
rule(s, Inches(5.05))
textbox(s, M, Inches(5.28), W - 2 * M, Inches(1.2),
        "Both remain valid for validation, so neither is switched off. RSA/SHA-1 has "
        "fallen from 81.7% of reverse DS records in 2009 to 4.9% today — the "
        "retirement is nearly complete, and this is the residue the pipeline's "
        "non-conformance signal now tracks.",
        13.5, INK2, spacing=1.3)

# ========================================================== 11. what we can't
s = slide()
header(s, "stated, not hidden", "What we can and cannot measure",
       "Every one of the 30 RFCs carries a verdict and the sentence that justifies "
       "it. An RFC we cannot answer is a documented result, not a silence.")
stat_row(s, [
    ("19", "measurable"),
    ("2", "partly measurable"),
    ("2", "ambiguous only"),
    ("7", "not measurable here"),
], y=Inches(2.50), color=INK)
rule(s, Inches(4.30))
bullets(s, [
    ("Resolver-side is invisible",
     "RFC 4035 and 8198 are properties of a query/response exchange. Authoritative "
     "measurement cannot see them — DNS-OARC's DITL could, but it is members-only."),
    ("Owner-name signatures",
     "RFC 9615 and 7672 are identified by a label (_signal, _25._tcp). Names are "
     "provenance in our signal model, not evidence."),
    ("A limit we can fix",
     "RFC 6840's mandatory-algorithm rules ARE visible, but need algorithm-set "
     "comparisons across records at one name. Our indicators are per-record."),
], y=Inches(4.55), step=Inches(0.86), label_w=Inches(3.6))

# ============================================================= 12. rigour
s = slide()
header(s, "what the checks caught", "Four corrections, before they reached a slide",
       "Reported because the process working is the point — each was found by a "
       "check rather than by someone noticing later.")
bullets(s, [
    ("A retracted headline",
     "\"Crossed 1%\" was wrong. A second composition break — RIPE reformatting in "
     "2015, its count falling 97% — was not being excluded. Corrected to 0.88-1.01%."),
    ("A publication-date leak",
     "A zero-scoring match kept its adoption verdict on a pre-publication "
     "observation. RFC 6840 showed a 2013 document first seen in 2010."),
    ("2.45M phantom matches",
     "RFC 6840 and 9364 matched \"any DNSSEC record\" — the evidence RFC 4033 "
     "already claims — for 25.9% of all matched rows. Both now match nothing."),
    ("An untested correlation",
     "\"Rollovers peak when adoption is steepest\" did not survive testing: r = 0.43, "
     "and the steepest year had the lowest rollover rate."),
], y=Inches(2.42), step=Inches(1.02), label_w=Inches(3.9))

# =============================================================== 13. next
s = slide()
header(s, "next", "Where this goes",
       "Measurement validity first; the corpus is now wide enough that breadth is "
       "no longer the binding constraint.")
bullets(s, [
    ("1  Run both corpora at full scale",
     "The nightly script is written and tested. Mirror once, scan locally, and the "
     "store stops being in the inner loop."),
    ("2  Cross-record indicators",
     "Unlocks RFC 6840's mandatory-algorithm rules — a real conformance check, and "
     "the kind of misconfiguration that actually breaks validation."),
    ("3  Splice the RIPE series",
     "Rather than dropping a whole RIR to avoid one 2015 format change. Needs "
     "someone who knows what the legacy zones held."),
    ("4  Ask about the closed sources",
     "DNS-OARC DITL is the only thing that would let us say anything about the "
     "resolver side. It needs a membership conversation, not a crawler."),
], y=Inches(2.42), step=Inches(1.02), label_w=Inches(3.9))
footnote(s, "Reproducible end to end: notebooks/dnssec_crossref_openintel_ripe.ipynb "
            "runs from a clean clone and re-derives every figure in this deck.")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

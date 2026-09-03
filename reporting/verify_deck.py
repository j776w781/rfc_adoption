"""Check every factual claim on the deck against the run data.

Exists because the deck accumulated stale figures faster than they could be
spotted by reading: a dek that contradicted its own table, residue numbers from
a superseded run, and a next-steps item asking for work that had already been
done. Each assertion below names the slide it defends.
"""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from pptx import Presentation

DECK = Path(sys.argv[1] if len(sys.argv) > 1
            else "out/analysis/dnssec_framework.pptx")
BUNDLE = Path(sys.argv[2] if len(sys.argv) > 2
              else "out/analysis/full_run_bundle.json")
PANEL = Path("out/panel_run/bottom_up.json")

bundle = json.loads(BUNDLE.read_text(encoding="utf-8"))
full = {r["label"]: r for r in bundle["bottom_up"]}
panel = {r["label"]: r for r in json.loads(PANEL.read_text(encoding="utf-8"))}
corpus = bundle["corpus"]

prs = Presentation(str(DECK))
slides = []
for sl in prs.slides:
    parts = []
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            parts.append(sh.text_frame.text)
        if sh.has_table:
            parts += [c.text for r in sh.table.rows for c in r.cells]
    slides.append(" ".join(parts))
text = " ".join(slides)

ok, bad = 0, []


def check(label, condition, detail=""):
    global ok
    if condition:
        ok += 1
    else:
        bad.append(f"{label}: {detail}")


def says(fragment, slide=None):
    hay = slides[slide - 1] if slide else text
    return fragment in hay


# --- corpus figures on the title slide ------------------------------------
tb = corpus["bytes"] / 1e12
check("title TB", says(f"{tb:.2f} TB", 1), f"deck should say {tb:.2f} TB")
check("title source-days", says(f"{corpus['source_days']:,}", 1),
      f"expected {corpus['source_days']:,}")
check("title files", says(f"{corpus['files']:,}", 1), f"expected {corpus['files']:,}")
check("title sources", says(f"{len(corpus['sources'])} sources", 1),
      f"expected {len(corpus['sources'])} sources")

# --- observable counts ----------------------------------------------------
observed = sum(1 for r in full.values() if r.get("t1_first_seen"))
check("observed count", says(f"14  →  {observed}") or says(f"14 → {observed}"),
      f"full run observed {observed}")

# --- Ed25519 / Ed448, the worked example ----------------------------------
ed, e4 = full["Ed25519"], full["Ed448"]
check("Ed25519 onset", says(f"{ed['onset_years']:.1f} y") or says(f"{ed['onset_years']:.1f}y"),
      f"expected {ed['onset_years']:.1f}")
check("Ed25519 first seen", says(ed["t1_first_seen"]), ed["t1_first_seen"])
check("Ed25519 partial", says(ed["t2_partial_usage"]), ed["t2_partial_usage"])
# Ed25519's shares are checked against the SEPARATED analysis further down, not
# the pooled bundle: a pooled share for it spans the 2024-01 break, and the deck
# deliberately quotes the per-corpus figures instead.
check("Ed448 onset", says(f"{e4['onset_years']:.1f} y"), f"expected {e4['onset_years']:.1f}")
check("Ed448 first seen", says(e4["t1_first_seen"]), e4["t1_first_seen"])
check("Ed25519 reached partial, Ed448 did not",
      ed["t2_partial_usage"] and not e4["t2_partial_usage"])

# --- the withdrawn monotonicity claim -------------------------------------
prim = ["ECC-GOST", "ECDSAP256SHA256", "ECDSAP384SHA384", "Ed25519", "Ed448"]
seq = [full[k]["onset_years"] for k in prim]
check("full-corpus sequence is not monotonic",
      not all(b >= a for a, b in zip(seq, seq[1:])), str(seq))
check("deck prints the sequence",
      says(" → ".join(f"{v:.1f}" for v in seq)),
      "expected " + " → ".join(f"{v:.1f}" for v in seq))
rev = [panel[k].get("onset_years") for k in prim if panel[k].get("onset_years")]
check("reverse sequence is monotonic", all(b >= a for a, b in zip(rev, rev[1:])), str(rev))

# --- residue figures ------------------------------------------------------
for lbl, key in (("SHA-1 signing", "SHA-1 signing still published"),
                 ("SHA-1 digest", "SHA-1 DS digest still published")):
    r = full[key]
    check(f"{lbl} share", says(f"{r['current_share_pct']:.2f}%"),
          f"expected {r['current_share_pct']:.2f}%")
    check(f"{lbl} records", says(f"{r['residue_records_after_publication']:,}"),
          f"expected {r['residue_records_after_publication']:,}")
check("ECC-GOST residue is zero", full["ECC-GOST still published"]["current_share_pct"] == 0)

# --- untestable observables ----------------------------------------------
nocov = [r for r in full.values() if r["state"] == "no_corpus_coverage"]
check("no_corpus_coverage count", says(f"{len(nocov)} observables"),
      f"expected {len(nocov)}")

# --- claims that must NOT appear -----------------------------------------
# The comparison slide legitimately prints the superseded reverse-only figures
# beside the corrected ones -- that is the point of it. Ban them everywhere else.
COMPARISON_SLIDE = next(i for i, t in enumerate(slides, 1)
                        if "What the full corpus changed" in t)
for frag in ("5.6 y", "5.8 y", "3.7 y", "4.0 y"):
    check(f"superseded {frag!r} appears only on the comparison slide",
          all(frag not in t for i, t in enumerate(slides, 1) if i != COMPARISON_SLIDE),
          f"{frag!r} found outside slide {COMPARISON_SLIDE}")

banned = {
    "the superseded 7.81% residue": "7.81%",
    "the superseded 21.77% residue": "21.77%",
    "'only one of them is used'": "Only one of them is used",
    "a request to run the full cache (already done)": "Run the pipeline over the full",
    "'Everything here is reverse-DNS only'": "Everything here is reverse-DNS only",
    "the old 2.5-5.8y primitive band": "2.5 – 5.8 y",
    "an uncaveated ECDSA peak-to-now pair": "76.2% 71.1%",
}
for why, frag in banned.items():
    check(f"does not claim {why}", not says(frag), f"found {frag!r}")

# --- the separated analysis, once the real timeline arrived ---------------
SPLIT = Path("out/server_run/split_analysis.json")
if SPLIT.exists():
    sp = json.loads(SPLIT.read_text(encoding="utf-8"))
    x = sp["cross_reference"]
    check("comparison month is where both corpora overlap",
          says(x["comparison_month"]), x["comparison_month"])
    comparable = [c for c in x["comparisons"] if c["difference_pct_points"] is not None]
    agree = [c for c in comparable if abs(c["difference_pct_points"]) <= 5]
    check("agreement count", says(f"{len(agree)} of {len(comparable)}"),
          f"expected {len(agree)} of {len(comparable)}")
    d2 = next(c for c in comparable if c["label"] == "SHA-256 DS digest")
    check("SHA-256 forward", says(f"{d2['forward_share_pct']:.2f}%"),
          f"{d2['forward_share_pct']:.2f}%")
    check("SHA-256 reverse", says(f"{d2['reverse_share_pct']:.2f}%"),
          f"{d2['reverse_share_pct']:.2f}%")
    fwd = {r["label"]: r for r in sp["sides"]["forward"]["bottom_up"]}
    rev = {r["label"]: r for r in sp["sides"]["reverse"]["bottom_up"]}
    check("Ed25519 forward onset", says(f"{fwd['Ed25519']['onset_years']:.1f} y"),
          str(fwd["Ed25519"]["onset_years"]))
    check("Ed25519 reverse onset", says(f"{rev['Ed25519']['onset_years']:.1f} y"),
          str(rev["Ed25519"]["onset_years"]))
    check("Ed25519 forward first seen", says(fwd["Ed25519"]["t1_first_seen"]))
    check("Ed25519 reverse first seen", says(rev["Ed25519"]["t1_first_seen"]))
    check("Ed25519 forward peak", says(f"{fwd['Ed25519']['peak_share_pct']:.2f}%"),
          f"{fwd['Ed25519']['peak_share_pct']:.2f}%")
    check("Ed25519 reverse peak", says(f"{rev['Ed25519']['peak_share_pct']:.2f}%"),
          f"{rev['Ed25519']['peak_share_pct']:.2f}%")
    check("no pooled Ed25519 peak is quoted", not says("2.03%"),
          "2.03% is the pooled figure and spans the break")
    # The 99.6% cliff and the bulk-migration figure are the two claims that
    # justify never pooling; both must be on the deck if the split is.
    check("the 99.6% cliff is stated", says("99.6%"))
    check("the bulk migration is stated", says("118,961"))
    check("no dangling cross-reference", "see note" not in text)

# --- internal consistency: every % on a slide must exist in the data ------
quoted = set(re.findall(r"\b(\d{1,3}\.\d{2})%", text))
known = set()
for r in full.values():
    for k in ("current_share_pct", "peak_share_pct", "residue_share_pct"):
        if r.get(k) is not None:
            known.add(f"{r[k]:.2f}")
for r in panel.values():
    for k in ("current_share_pct", "peak_share_pct"):
        if r.get(k) is not None:
            known.add(f"{r[k]:.2f}")
if SPLIT.exists():
    for side in sp["sides"].values():
        for r in side["bottom_up"]:
            for k in ("current_share_pct", "peak_share_pct"):
                if r.get(k) is not None:
                    known.add(f"{r[k]:.2f}")
    for c in sp["cross_reference"]["comparisons"]:
        for k in ("forward_share_pct", "reverse_share_pct"):
            if c.get(k) is not None:
                known.add(f"{c[k]:.2f}")
known |= {"74.60", "22.30", "3.10", "0.02",
          "2.04", "0.00", "0.55", "0.003", "99.60", "2.03"}   # Rogers illustration + panel prose
unknown = sorted(quoted - known)
check("every two-decimal percentage traces to the data", not unknown,
      f"untraceable: {unknown}")

print(f"{ok} checks passed, {len(bad)} failed")
for b in bad:
    print(f"  FAIL  {b}")
sys.exit(1 if bad else 0)
